#!/usr/bin/env python3
"""Build the ADM EBR interview deck as PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Colors
BG = RGBColor(0x0B, 0x0F, 0x14)
SURFACE = RGBColor(0x18, 0x20, 0x2C)
TEXT = RGBColor(0xE8, 0xED, 0xF5)
MUTED = RGBColor(0x9A, 0xA8, 0xBC)
ACCENT = RGBColor(0x6E, 0xE7, 0xB7)
ACCENT2 = RGBColor(0x60, 0xA5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x22, 0x30)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_kicker(slide, text, top=Inches(0.45)):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Segoe UI"
    return box


def add_title(slide, text, top=Inches(0.85), size=32):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.font.name = "Segoe UI"
    return box


def add_body(slide, lines, left=Inches(0.6), top=Inches(1.9), width=Inches(12.1), size=16, color=MUTED):
    height = Inches(min(5.2, 0.32 * max(len(lines), 1) + 0.5))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
    return box


def add_bullets(slide, items, left=Inches(0.6), top=Inches(1.75), width=Inches(12.1), size=17, color=TEXT, level=0):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.level = level
        p.space_after = Pt(10)
    return box


def add_slide_number(slide, num, total=18):
    box = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = f"{num:02d}"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.RIGHT


def add_table(slide, headers, rows, left=Inches(0.5), top=Inches(1.65), width=Inches(12.3)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(min(5.0, 0.38 * n_rows + 0.5)))
    table = table_shape.table

    col_widths = [width / n_cols] * n_cols
    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = MUTED
        p.font.name = "Segoe UI"

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x20) if r % 2 else BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT
            p.font.name = "Segoe UI"
            cell.text_frame.word_wrap = True
            cell.vertical_anchor = MSO_ANCHOR.TOP
    return table_shape


def add_two_column_cards(slide, cards, top=Inches(1.7)):
    col_w = Inches(5.95)
    gap = Inches(0.4)
    for i, (title, bullets) in enumerate(cards):
        left = Inches(0.6) + i * (col_w + gap)
        shape = slide.shapes.add_shape(1, left, top, col_w, Inches(5.0))  # rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = RGBColor(0x2A, 0x35, 0x48)

        tbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), col_w - Inches(0.4), Inches(0.5))
        p = tbox.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT2
        p.font.name = "Segoe UI"

        bbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), col_w - Inches(0.4), Inches(4.2))
        tf = bbox.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(12)
            p.font.color.rgb = MUTED
            p.font.name = "Segoe UI"
            p.space_after = Pt(6)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    n = 0

    def slide():
        nonlocal n
        n += 1
        s = prs.slides.add_slide(blank)
        set_slide_bg(s)
        return s, n

    # 01 Title
    s, num = slide()
    add_kicker(s, "AI Deployment Manager — Interview Exercise")
    add_title(s, "Executive Business Review", top=Inches(2.8), size=44)
    add_body(s, [
        "Portfolio prioritization, champion strategy, and a live EBR role-play with Uber Engineering Leadership.",
        "",
        "Gabriel Smith  ·  Confidential — fictional scenario data  ·  July 2026",
    ], top=Inches(4.2), size=18)
    add_slide_number(s, num)

    # 02 Agenda
    s, num = slide()
    add_kicker(s, "Overview")
    add_title(s, "Agenda", size=36)
    add_bullets(s, [
        "Portfolio prioritization — top 5 accounts & 2 deprioritized",
        "Champion strategy — 5 individuals at Uber (prioritized account)",
        "EBR role-play — Business review with Uber Engineering",
        "Product roadmap — Cloud Agents, Automations, Bugbot + live demo",
        "Strategic initiatives — where to prioritize & 90-day next steps",
        "Asks for Cursor leadership",
    ], top=Inches(1.85), size=20)
    add_slide_number(s, num)

    # 03 Framework
    s, num = slide()
    add_kicker(s, "Portfolio")
    add_title(s, "Prioritization framework", size=34)
    cards = [
        ("ADM fit", ["Metrics & ROI storytelling", "Enterprise rollouts", "Champion building"]),
        ("Scale & whitespace", ["4K+ SWEs", "<60% active or <15% power-user penetration"]),
        ("Retention urgency", ["Health & renewal timing", "ARR at risk", "Competitor trends"]),
        ("Growth trajectory", ["Positive/recoverable MoM active trend", "Limited competitor displacement"]),
        ("Strategic value", ["Partnerships & references", "Ecosystem unlock (e.g. Databricks)"]),
        ("Expansion runway", ["Low cloud/agent usage", "Automations & Bugbot upsell"]),
    ]
    top = Inches(1.65)
    for row in range(2):
        row_cards = cards[row * 3:(row + 1) * 3]
        for i, (title, bullets) in enumerate(row_cards):
            left = Inches(0.55) + i * Inches(4.15)
            shape = s.shapes.add_shape(1, left, top + row * Inches(2.55), Inches(3.95), Inches(2.35))
            shape.fill.solid()
            shape.fill.fore_color.rgb = SURFACE
            shape.line.color.rgb = RGBColor(0x2A, 0x35, 0x48)
            tbox = s.shapes.add_textbox(left + Inches(0.15), top + row * Inches(2.55) + Inches(0.12), Inches(3.65), Inches(2.1))
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = ACCENT2
            p.font.name = "Segoe UI"
            for b in bullets:
                p = tf.add_paragraph()
                p.text = f"• {b}"
                p.font.size = Pt(11)
                p.font.color.rgb = MUTED
                p.font.name = "Segoe UI"
    add_body(s, ["12-account mixed book · $195K–$1.4M ARR · renewal 3–9 months"], top=Inches(6.85), size=11)
    add_slide_number(s, num)

    # 04 Top 5
    s, num = slide()
    add_kicker(s, "Portfolio")
    add_title(s, "Top 5 accounts to prioritize (next 90 days)", size=28)
    add_table(s,
        ["Account", "ARR", "Health", "SWEs", "Active%", "Cloud%", "Why prioritize"],
        [
            ["LinkedIn", "$1.40M", "Red", "10,025", "25%", "3.7%", "Highest ARR at risk; hackathon interest; massive whitespace; +31% competitor. #1 priority."],
            ["Uber", "$980K", "Green", "15,423", "54%", "10.3%", "Post-rollout gap; expand into 2 eng orgs; build champions. Largest SWE base."],
            ["Spotify", "$1.00M", "Yellow", "3,525", "40%", "10.2%", "Exec wants expansion but needs adoption-depth proof. Retention vs. +39% competitor."],
            ["Databricks", "$420K", "Yellow", "5,480", "80%", "19.2%", "Strategic partnership; under-monetized; QBR this quarter; competitor declining."],
            ["Workday", "$1.25M", "Yellow", "7,095", "93%", "33.1%", "+50% MoM growth; ADM runs consumption analysis for renewal ROI story."],
        ],
        top=Inches(1.55))
    add_body(s, ["Note: Strong alternate for Workday is Stripe ($1.35M, Red, 4 mo renewal, 2.5% cloud)."], top=Inches(6.7), size=10)
    add_slide_number(s, num)

    # 05 Deprioritize
    s, num = slide()
    add_kicker(s, "Portfolio")
    add_title(s, "2 accounts to deprioritize", size=34)
    add_two_column_cards(s, [
        ("Figma — $280K · 737 SWEs · 3 mo renewal", [
            "100% penetration but declining: −12% MoM active, −15.75% competitor",
            "Active competitor pilot — ADM can't add headcount",
            "Competitive defense / product fit, not rollout work",
            "Already end-to-end; limited 90-day ADM leverage",
        ]),
        ("Notion — $195K · 874 SWEs · 3 mo renewal", [
            "Stuck on unresolved product/security concerns, not growth",
            "ADM can escalate feedback but can't fix blockers in 90 days",
            "Below 4K SWE bar; small commercial footprint",
            "Oversold position damages trust — route to Product/Eng first",
        ]),
    ])
    add_body(s, ["Also on watch: Disney (12K SWE, no exec sponsor) · Rivian (AE deprioritized) · Gusto (below 4K SWE)"], top=Inches(6.75), size=11)
    add_slide_number(s, num)

    # 06 Champions intro
    s, num = slide()
    add_kicker(s, "Champions")
    add_title(s, "Champion strategy — Uber", size=34)
    add_body(s, [
        "Selected account: largest SWE footprint, post-rollout whitespace, AE-identified expansion into two additional eng orgs.",
    ], top=Inches(1.55), size=16, color=TEXT)
    add_two_column_cards(s, [
        ("Goal", ["Convert power users into org-wide advocates", "Establish rollout owner where none exists"]),
        ("Signal", ["Strong Q-1 onboarding, then engagement drop", "Classic champion vacuum after rollout"]),
        ("Outcome", ["3 champions activated in 90 days", "Exec EBR with adoption-depth metrics & expansion case"]),
    ], top=Inches(2.2))
    # Fix 3-col on champions intro - use manual layout
    add_slide_number(s, num)

    # 07 Champions table - split across content
    s, num = slide()
    add_kicker(s, "Champions")
    add_title(s, "5 champion targets at Uber (1/2)", size=30)
    add_table(s,
        ["Name", "Role", "Why champion", "First meeting (30 min)"],
        [
            ["Priya Sharma", "Dir Eng, Developer Platform", "Owns dev tooling standards; natural buyer for org-wide rollout governance.", "Listen: rollout gaps. Show: peer benchmark (54% vs 70%+). Ask: co-own 90-day depth playbook."],
            ["Marcus Webb", "Sr EM, Mobility Core", "Manages expansion org; 92% agent usage, 10% cloud — ready for automation upsell.", "Listen: team usage today. Show: power-user teardown. Ask: pilot Automations for PR review."],
            ["Elena Vasquez", "Staff Eng, AI DevEx", "Internal AI tooling advocate; de facto power user; scales best practices.", "Listen: post-Q1 blockers. Show: Cloud Agent + Automations demo. Ask: co-host Agentic Dev lunch & learn."],
        ],
        top=Inches(1.55))
    add_slide_number(s, num)

    s, num = slide()
    add_kicker(s, "Champions")
    add_title(s, "5 champion targets at Uber (2/2)", size=30)
    add_table(s,
        ["Name", "Role", "Why champion", "First meeting (30 min)"],
        [
            ["James Okonkwo", "VP Eng, Marketplace", "Exec sponsor for 2nd expansion org; needs ROI proof before commercial conversation.", "Listen: expansion criteria. Show: 90-day adoption dashboard. Ask: quarterly EBR cadence."],
            ["Sofia Lindström", "Principal Eng, Infrastructure", "Credible technical voice; influences security reviews; unblocks admin bottleneck.", "Listen: compliance concerns. Show: enterprise controls + Bugbot rules. Ask: join admin working group."],
        ],
        top=Inches(1.55))
    add_slide_number(s, num)

    # 08 EBR intro
    s, num = slide()
    add_kicker(s, "EBR Role Play")
    add_title(s, "Executive Business Review — Uber", size=32)
    add_body(s, [
        "Simulated QBR with James Okonkwo (VP Eng, Marketplace) + Priya Sharma (Dir, Developer Platform)",
        "Account: $980K ARR · 15,423 SWEs · Green health · 7 mo to renewal",
    ], top=Inches(1.55), size=17, color=TEXT)
    add_two_column_cards(s, [
        ("Attendees — Cursor", [
            "ADM — Gabriel Smith (facilitator)",
            "Account Executive — partnership & commercial",
            "Solutions Engineer — technical deep-dives",
        ]),
        ("Attendees — Uber", [
            "James Okonkwo — VP Eng, Marketplace",
            "Priya Sharma — Dir, Developer Platform",
            "Elena Vasquez — Staff Eng, AI DevEx (optional)",
        ]),
    ], top=Inches(2.5))
    add_slide_number(s, num)

    # 09 Usage
    s, num = slide()
    add_kicker(s, "EBR · Business Review")
    add_title(s, "Usage & partnership to date", size=32)
    add_two_column_cards(s, [
        ("Adoption snapshot", [
            "L30D active users: 8,385 (54% of SWEs)",
            "Power users: 2,757 (17.9% of SWEs)",
            "MoM active trend: +15%",
            "Agent usage: 92% · Composer: 39.9%",
            "Cloud Agent usage: 10.3%",
            "Competitor trend: +7.2% MoM",
        ]),
        ("Partnership timeline", [
            "Q-1: Enterprise onboarding — strong rollout across Mobility & Platform",
            "Q0 (now): Engagement plateau; no dedicated rollout owner",
            "AE pipeline: Two additional eng orgs; admin slow on intros",
            "Renewal: 7 months — healthy sentiment; depth proof needed for expansion",
        ]),
    ])
    add_slide_number(s, num)

    # 10 Value
    s, num = slide()
    add_kicker(s, "EBR · Business Review")
    add_title(s, "Value realized", size=34)
    cards = [
        ("Velocity", "~2.4 hrs saved per dev/week on boilerplate, tests, refactors (survey n=120)"),
        ("Quality", "18% faster PR cycle time in Mobility Core (Q-1 cohort)"),
        ("Scale", "8,385 engineers active in 90 days — fastest enterprise ramp in book"),
    ]
    top = Inches(1.65)
    for i, (title, desc) in enumerate(cards):
        left = Inches(0.55) + i * Inches(4.15)
        shape = s.shapes.add_shape(1, left, top, Inches(3.95), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = RGBColor(0x2A, 0x35, 0x48)
        tbox = s.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), Inches(3.65), Inches(1.8))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT2
        p.font.name = "Segoe UI"
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED
        p.font.name = "Segoe UI"
    add_body(s, [
        "ROI framing: 2.4 hrs/dev/week × 8,385 active × $85/hr ≈ $1.7M/month recovered capacity vs. $980K annual contract.",
        "Expansion to remaining 46% of SWEs + two new orgs could 2–3× this impact.",
    ], top=Inches(4.0), size=15, color=TEXT)
    add_slide_number(s, num)

    # 11 Strategic initiatives
    s, num = slide()
    add_kicker(s, "EBR · Business Review")
    add_title(s, "Strategic initiatives supported", size=30)
    add_table(s,
        ["Uber initiative", "How Cursor supports it", "Status"],
        [
            ["Platform modernization", "Agent-assisted migrations; Composer for large-scale refactors", "Active"],
            ["AI-native developer experience", "Cloud Agents for async tasks; Automations for CI/CD workflows", "Early — 10% cloud"],
            ["Engineering efficiency metrics", "ADM adoption dashboard; power-user cohort analysis", "In progress"],
            ["Security & code quality", "Bugbot on PRs; custom rules for service mesh patterns", "Pilot — 2 squads"],
            ["Cross-org eng standards", "Shared .cursor/rules + champion guild (Elena leading)", "Planned Q+1"],
        ],
        top=Inches(1.55))
    add_slide_number(s, num)

    # 12 Feedback
    s, num = slide()
    add_kicker(s, "EBR · Business Review")
    add_title(s, "Partnership feedback", size=34)
    add_two_column_cards(s, [
        ("Highlights", [
            "Fastest enterprise onboarding — 54% active in 90 days",
            "Strong agent adoption (92%) — daily driver, not a toy",
            "Responsive AE + ADM cadence during rollout",
            "Power-user Slack community (400+ members)",
        ]),
        ("Areas for improvement", [
            "Post-rollout drop-off — need quarterly touchpoints",
            "Admin bottleneck on cross-org introductions",
            "Cloud Agent awareness: only 10% using cloud",
            "Leadership wants team-level dashboards",
            "Competitor noise: +7% MoM in one Mobility squad",
        ]),
    ])
    add_slide_number(s, num)

    # 13 Roadmap
    s, num = slide()
    add_kicker(s, "EBR · Product Roadmap")
    add_title(s, "What's next — Cursor product roadmap", size=28)
    cards = [
        ("Cloud Agents", [
            "Isolated VMs with full dev environments",
            "Parallel agents; MCP, multi-repo, artifacts",
            "Trigger: IDE, web, Slack, GitHub, Linear, API",
            "Uber fit: async refactors across monorepo",
        ]),
        ("Automations", [
            "Schedule or event-triggered agents",
            "PR, Slack, PagerDuty, webhooks",
            "Marketplace templates + memory",
            "Uber fit: PR triage, on-call → fix PR",
        ]),
        ("Bugbot", [
            "Auto-review on every PR; custom rules",
            "70%+ flags resolved before merge",
            "Autofix spawns Cloud Agent to fix & push",
            "Uber fit: enforce service-mesh patterns",
        ]),
    ]
    top = Inches(1.6)
    for i, (title, bullets) in enumerate(cards):
        left = Inches(0.55) + i * Inches(4.15)
        shape = s.shapes.add_shape(1, left, top, Inches(3.95), Inches(5.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = RGBColor(0x2A, 0x35, 0x48)
        tbox = s.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), Inches(3.65), Inches(4.7))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT
        p.font.name = "Segoe UI"
        for b in bullets:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(11)
            p.font.color.rgb = MUTED
            p.font.name = "Segoe UI"
    add_slide_number(s, num)

    # 14 Demo
    s, num = slide()
    add_kicker(s, "EBR · Demo")
    add_title(s, "Demo: Automations with Cloud Agents", size=30)
    add_body(s, [
        "Weekly PR quality automation for Uber Mobility Core",
        "",
        "Trigger: Every Monday 9am PT · GitHub PR merged to mobility-core",
        "",
        "1. Cloud Agent clones repo, reviews last 7 days of merged PRs",
        "2. Runs targeted test suite in isolated VM; captures screenshot artifacts",
        "3. Posts summary to #mobility-eng Slack with top 3 findings",
        "4. If critical: opens draft PR with Autofix via Bugbot integration",
        "",
        "90-day pilot: Mobility Core (1 squad) → Marketplace → org-wide Automations library",
    ], top=Inches(1.55), size=15, color=TEXT)
    add_slide_number(s, num)

    # 15 Next steps
    s, num = slide()
    add_kicker(s, "EBR · Strategic Initiatives")
    add_title(s, "Where to prioritize & next steps", size=28)
    add_table(s,
        ["Priority", "Initiative", "Owner", "Timeline", "Success metric"],
        [
            ["P0", "Appoint rollout owner + unblock 2 eng org intros", "Priya + ADM", "Wk 1–2", "Intros scheduled"],
            ["P0", "Adoption-depth dashboard for exec visibility", "ADM", "Wk 2–4", "Team-level active/cloud %"],
            ["P1", "Automations pilot (Mobility Core PR review)", "Marcus + Elena", "Wk 4–8", "Cloud 10% → 20%"],
            ["P1", "Bugbot rollout to 5 squads with custom rules", "Sofia + SE", "Wk 6–10", "70%+ flag resolution"],
            ["P2", "Quarterly EBR + expansion business case", "James + AE", "Wk 10–12", "Expansion proposal for 2 orgs"],
        ],
        top=Inches(1.55))
    add_slide_number(s, num)

    # 16 Asks 1-3
    s, num = slide()
    add_kicker(s, "Asks for Cursor Leadership")
    add_title(s, "3 areas where I need leadership support", size=28)
    add_bullets(s, [
        "1. Ongoing product ↔ field forum sessions",
        "   Ask: Bi-weekly ADM + Product + Eng forum for field signal and roadmap visibility.",
        "   Question: What is our forum today? Is there a structured channel, or do ADMs route feedback ad hoc?",
        "",
        "2. Cursor exec leadership to unlock next-level partnerships",
        "   Ask: Exec-to-exec alignment for Databricks, LinkedIn, Uber — not just ADM + VP Eng.",
        "   Why: Accelerates procurement, references, and co-marketing on strategic accounts.",
        "",
        "3. Quarterly feedback sessions with leadership",
        "   Ask: 30-min quarterly 1:1 for coaching on portfolio strategy and career growth.",
        "   Question: Where am I falling short? Am I overweighting expansion vs. retention?",
    ], top=Inches(1.55), size=14)
    add_slide_number(s, num)

    # 17 Alignment
    s, num = slide()
    add_kicker(s, "Asks for Cursor Leadership")
    add_title(s, "Alignment: which accounts are actually strategic?", size=26)
    add_two_column_cards(s, [
        ("The tension", [
            "12 accounts; framework says whitespace + ADM fit.",
            "Leadership may have different bets (Databricks, Workday, Disney).",
            "",
            "Question: How much of my time should be saving vs. expanding?",
            "• Saving: LinkedIn, Stripe, Spotify",
            "• Expanding: Uber, Databricks",
            "• Monitoring: Workday (healthy, limited ADM ROI)",
        ]),
        ("Proposed 90-day time allocation", [
            "Retention / save — 40%",
            "  LinkedIn, Spotify, Stripe check-ins",
            "",
            "Expansion / depth — 40%",
            "  Uber, Databricks, champion programs",
            "",
            "Strategic / light touch — 20%",
            "  Workday consumption, Disney watch",
            "",
            "Seeking leadership sign-off on allocation.",
        ]),
    ])
    add_slide_number(s, num)

    # 18 Close
    s, num = slide()
    add_kicker(s, "Thank you")
    add_title(s, "Questions & discussion", top=Inches(2.8), size=44)
    add_body(s, [
        "Ready to role-play the Uber EBR live, walk through the Automations demo,",
        "or deep-dive any account in the portfolio.",
        "",
        "Gabriel Smith · ADM Interview · Portfolio exercise data confidential",
    ], top=Inches(4.2), size=18)
    add_slide_number(s, num)

    out = Path(__file__).parent / "Executive_Business_Review_ADM_Interview.pptx"
    prs.save(str(out))
    print(f"Saved: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
